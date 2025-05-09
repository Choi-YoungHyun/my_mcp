from typing import List, Optional, Any
import os
from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_chroma import Chroma

from rag.base import RetrievalChain

class PPTRetrievalChain(RetrievalChain):
    """
    PPT-specific implementation of the RetrievalChain.
    
    This class specializes in loading, splitting, and indexing PPT documents
    for retrieval. 
    """
    
    def __init__(self, 
                    source_uri: List[str], 
                    persist_directory: Optional[str] = None,
                    **kwargs) -> None:
        """
        Initialize a PPT retrieval chain.
        
        Args:
            source_uri: List of PPT file paths
            persist_directory: Directory to persist vector store
            **kwargs: Additional keyword arguments for the base RetrievalChain
        """

        super().__init__(source_uri=source_uri, persist_directory=persist_directory, **kwargs)
    
    def load_documents(self, source_uris: List[str]) -> List[Document]:
        print(">>>>>>>>>>>>>")
        print(source_uris)
        print(">>>>>>>>>>>>>")
        """
        Load PPT documents from file paths.
        
        Args:
            source_uris: List of PPT file paths
            
        Returns:
            List of loaded documents
        """

        docs = []
        for source_uri in source_uris:
            if not os.path.exists(source_uri):
                print(f"File not found: {source_uri}")
                continue
                
            print(f"Loading PPT: {source_uri}")
            presentation = Presentation(source_uri)
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                for shape in slide.shapes:
                    #객체에 text 속성이 있는지 없는지
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    doc = Document(
                        page_content="\n".join(slide_text),
                        metadata={"source": source_uri, "slide": i}
                    )
                    docs.append(doc)
        
        return docs
    
    def create_text_splitter(self) -> RecursiveCharacterTextSplitter:
        """
        Create a text splitter optimized for PPT documents.
        
        Returns:
            A text splitter instance suitable for PPTs
        """
        
        return RecursiveCharacterTextSplitter(
            chunk_size=600,
            chunk_overlap=50
        )
    
    def create_vectorstore(self, split_docs: List[Document]) -> Any:
        """
        Create a vector store from split PPT documents.
        
        Args:
            split_docs: Split document chunks
            
        Returns:
            A vector store instance
            
        Raises:
            ValueError: If there are no split documents
        """
        
        if not split_docs:
            raise ValueError("No split documents available.")
            
        if self.persist_directory:
            os.makedirs(self.persist_directory, exist_ok=True)
            
            if os.path.exists(self.persist_directory) and any(os.listdir(self.persist_directory)):
                print(f"Loading existing vector store: {self.persist_directory}")

                return Chroma(
                    persist_directory=self.persist_directory,
                    embedding_function=self.create_embedding()
                )
        
        print("Creating new vector store...")

        vectorstore = Chroma.from_documents(
            documents=split_docs,
            embedding=self.create_embedding(),
            persist_directory=self.persist_directory
        )
        
        return vectorstore